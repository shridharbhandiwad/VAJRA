#!/usr/bin/env python3
"""
Script to create a comprehensive PowerPoint presentation for the Radar System Monitoring Application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(44)
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    
    # Add content
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        if isinstance(item, dict):
            p.text = item['text']
            p.level = item.get('level', 0)
        else:
            p.text = item
            p.level = 0
        
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    return slide

def add_diagram_slide(prs, title, diagram_text):
    """Add a slide with diagram/code content"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add diagram box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
    shape.line.color.rgb = RGBColor(0, 102, 204)
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = diagram_text
    p.font.name = 'Courier New'
    p.font.size = Pt(10)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Radar System Monitoring Application",
        "Architecture, Design, Patterns & Implementation\nComprehensive Technical Presentation"
    )
    
    # Slide 2: Table of Contents
    add_content_slide(prs, "Table of Contents", [
        "1. Project Overview",
        "2. System Architecture",
        "3. Design Principles & Patterns",
        "4. Core Components",
        "5. Implementation Details",
        "6. Key Features",
        "7. Block Diagram",
        "8. UML Class Diagram",
        "9. Sequence Diagrams",
        "10. Deployment Architecture",
        "11. Technology Stack",
        "12. Conclusion"
    ])
    
    # Slide 3: Project Overview
    add_content_slide(prs, "Project Overview", [
        "• Comprehensive Qt-based radar system monitoring application",
        "• Real-time health monitoring and visualization",
        "• Modular architecture with zero-code component extension",
        "• Designer and Runtime modes with role-based access",
        "• Multi-protocol support (TCP, UDP, WebSocket, MQTT)",
        "• Advanced analytics dashboard with 10 chart types",
        "• External Python simulators for testing"
    ])
    
    # Slide 4: System Architecture - High Level
    add_content_slide(prs, "System Architecture - High Level", [
        "Three-tier architecture:",
        {"text": "Presentation Layer - Qt GUI", "level": 1},
        {"text": "Designer Mode (visual layout creation)", "level": 2},
        {"text": "Runtime Mode (real-time monitoring)", "level": 2},
        {"text": "Analytics Dashboard (data visualization)", "level": 2},
        {"text": "Business Logic Layer - Core Components", "level": 1},
        {"text": "Component Registry (modular management)", "level": 2},
        {"text": "Message Server (health data processing)", "level": 2},
        {"text": "Analytics Engine (metrics calculation)", "level": 2},
        {"text": "Data Layer", "level": 1},
        {"text": "JSON configuration (components.json)", "level": 2},
        {"text": "Design files (.design format)", "level": 2},
        {"text": "Real-time health data streams", "level": 2}
    ])
    
    # Slide 5: Architecture Diagram
    arch_diagram = """
┌─────────────────────────────────────────────────────────┐
│                 Unified Application (Qt)                 │
├─────────────────────────────────────────────────────────┤
│  ┌──────────────┐  ┌──────────────┐  ┌──────────────┐  │
│  │   Designer   │  │   Runtime    │  │  Analytics   │  │
│  │     Mode     │  │    Monitor   │  │  Dashboard   │  │
│  └──────┬───────┘  └──────┬───────┘  └──────┬───────┘  │
│         │                  │                  │          │
│  ┌──────┴──────────────────┴──────────────────┴───────┐ │
│  │            Component Registry (Singleton)           │ │
│  │  Loads from components.json - Zero-code Extension   │ │
│  └──────────┬──────────────┬──────────────┬────────────┘ │
│             │              │              │               │
│    ┌────────┴────┐  ┌──────┴──────┐  ┌──────┴──────┐    │
│    │  Component  │  │   Canvas    │  │  Message    │    │
│    │    List     │  │  (Designer) │  │   Server    │    │
│    └─────────────┘  └─────────────┘  └──────┬──────┘    │
│                                              │           │
└──────────────────────────────────────────────┼───────────┘
                                               │
                     ┌─────────────────────────┴─────┐
                     │  TCP (12345) / UDP (12346)    │
                     └─────────────────┬─────────────┘
                                       │
          ┌────────────────────────────┴──────────────────┐
          │    External Subsystem Simulators (Python)     │
          │  ┌──────┐  ┌──────┐  ┌──────┐  ┌──────┐      │
          │  │comp_1│  │comp_2│  │comp_3│  │comp_N│      │
          │  └──────┘  └──────┘  └──────┘  └──────┘      │
          └───────────────────────────────────────────────┘
"""
    add_diagram_slide(prs, "System Architecture Diagram", arch_diagram)
    
    # Slide 6: Design Principles
    add_content_slide(prs, "Software Design Principles", [
        "• Single Responsibility Principle",
        {"text": "Each class has one well-defined purpose", "level": 1},
        "• Open/Closed Principle",
        {"text": "Open for extension, closed for modification", "level": 1},
        {"text": "New components added via JSON without code changes", "level": 1},
        "• Dependency Inversion",
        {"text": "Depend on abstractions (ComponentRegistry)", "level": 1},
        "• Separation of Concerns",
        {"text": "UI, Business Logic, and Data layers separated", "level": 1},
        "• DRY (Don't Repeat Yourself)",
        {"text": "Data-driven rendering eliminates code duplication", "level": 1}
    ])
    
    # Slide 7: Design Patterns
    add_content_slide(prs, "Design Patterns Implemented", [
        "• Singleton Pattern",
        {"text": "ComponentRegistry - single instance manages all types", "level": 1},
        "• Registry Pattern",
        {"text": "JSON-based component type registration", "level": 1},
        "• Observer Pattern",
        {"text": "MessageServer notifies components of health updates", "level": 1},
        "• Factory Pattern",
        {"text": "Component creation based on registry definitions", "level": 1},
        "• Model-View-Controller (MVC)",
        {"text": "Separation of data, presentation, and control logic", "level": 1},
        "• Strategy Pattern",
        {"text": "Multi-protocol support (TCP, UDP, WebSocket, MQTT)", "level": 1}
    ])
    
    # Slide 8: Core Components Overview
    add_two_column_slide(prs, "Core Components", 
        [
            "• ComponentRegistry",
            "  - Singleton managing component types",
            "  - Loads from components.json",
            "  - Validates definitions",
            "",
            "• Component",
            "  - Graphics item for canvas",
            "  - Data-driven rendering",
            "  - Health status visualization",
            "",
            "• Canvas",
            "  - Designer drag-drop area",
            "  - Connection management",
            "  - Sub-component support",
            "",
            "• MessageServer",
            "  - Multi-protocol receiver",
            "  - TCP & UDP listeners",
            "  - Health data dispatcher"
        ],
        [
            "• MainWindow",
            "  - Central application window",
            "  - Mode switching (Designer/Runtime)",
            "  - Toolbar and menu management",
            "",
            "• AnalyticsDashboard",
            "  - 10 chart types",
            "  - 6 KPI metrics",
            "  - Real-time updates",
            "",
            "• LoginDialog",
            "  - Role-based authentication",
            "  - User: Designer / User",
            "",
            "• ThemeManager",
            "  - Dark/Light theme switching",
            "  - Persistent preferences"
        ]
    )
    
    # Slide 9: Implementation - Component Registry
    add_content_slide(prs, "Implementation: Component Registry", [
        "• Singleton class managing component definitions",
        "• Key methods:",
        {"text": "loadFromJson() - Parse components.json", "level": 1},
        {"text": "addComponentType() - Runtime type addition", "level": 1},
        {"text": "getDefinition() - Retrieve component info", "level": 1},
        {"text": "getAllTypes() - List all registered types", "level": 1},
        "• ComponentDefinition structure:",
        {"text": "type_id, display_name, label, description", "level": 1},
        {"text": "image_dir, icon_color, shape", "level": 1},
        {"text": "subsystems[], protocol, port, category", "level": 1},
        "• Benefits:",
        {"text": "Zero code changes for new component types", "level": 1},
        {"text": "Backward compatibility with existing designs", "level": 1}
    ])
    
    # Slide 10: Implementation - Message Server
    add_content_slide(prs, "Implementation: Message Server", [
        "• Multi-protocol health data receiver",
        "• Protocols supported:",
        {"text": "TCP - Port 12345 (reliable, line-delimited JSON)", "level": 1},
        {"text": "UDP - Port 12346 (low-latency, datagram JSON)", "level": 1},
        "• Message format:",
        {"text": "{ \"component_id\": \"comp_1\", \"color\": \"#00FF00\", \"size\": 95.5 }", "level": 1},
        "• Processing flow:",
        {"text": "1. Receive JSON message", "level": 1},
        {"text": "2. Parse and validate", "level": 1},
        {"text": "3. Find component on canvas", "level": 1},
        {"text": "4. Update visual indicators", "level": 1},
        {"text": "5. Trigger analytics update", "level": 1},
        "• Multi-client support with concurrent connections"
    ])
    
    # Slide 11: Key Features - Designer Mode
    add_content_slide(prs, "Key Features: Designer Mode", [
        "• Visual component placement",
        {"text": "Drag-and-drop from component panel", "level": 1},
        {"text": "Component resizing and positioning", "level": 1},
        "• Connection drawing",
        {"text": "Uni-directional and bi-directional connections", "level": 1},
        {"text": "Custom labels and colors", "level": 1},
        "• Sub-component support",
        {"text": "Labels, LineEdits, Buttons inside components", "level": 1},
        {"text": "Drag-drop with auto-containment", "level": 1},
        "• Add component types at runtime",
        {"text": "\"+ Add Component Type\" button", "level": 1},
        {"text": "No code changes required", "level": 1},
        "• Save/Load design files (.design format)"
    ])
    
    # Slide 12: Key Features - Runtime Mode
    add_content_slide(prs, "Key Features: Runtime Mode", [
        "• Real-time health monitoring",
        {"text": "Color-coded status indicators", "level": 1},
        {"text": "Green (90-100%) → Yellow (70-89%) → Orange (40-69%)", "level": 1},
        {"text": "Red (10-39%) → Gray (0-9% offline)", "level": 1},
        "• Voice alerts for critical states",
        "• Enlarged component views",
        {"text": "Health trend charts", "level": 1},
        {"text": "Subsystem details", "level": 1},
        "• Analytics tracking",
        {"text": "Health updates count", "level": 1},
        {"text": "Status changes", "level": 1},
        {"text": "Health level transitions", "level": 1},
        "• Multi-protocol message reception (TCP & UDP)"
    ])
    
    # Slide 13: Key Features - Analytics Dashboard
    add_content_slide(prs, "Key Features: Analytics Dashboard", [
        "• 10 different chart types:",
        {"text": "1. Health Trend (Multi-line Spline)", "level": 1},
        {"text": "2. Component Distribution (Donut Pie)", "level": 1},
        {"text": "3. Subsystem Performance (Vertical Bar)", "level": 1},
        {"text": "4. Health Distribution (Stacked Area)", "level": 1},
        {"text": "5. Message Frequency (Scatter)", "level": 1},
        {"text": "6. Telemetry Streams (Multi-line)", "level": 1},
        {"text": "7. Alert History (Grouped Bar)", "level": 1},
        {"text": "8. Component Comparison (Horizontal Bar)", "level": 1},
        {"text": "9. System Efficiency (Gradient Area)", "level": 1},
        {"text": "10. Uptime Analysis (Stacked Bar)", "level": 1},
        "• 6 KPI cards: Total/Active Components, Avg Health, Alerts, Efficiency, Uptime",
        "• Auto-refresh every 5 seconds",
        "• Dark/Light theme support"
    ])
    
    # Slide 14: Block Diagram
    block_diagram = """
┌──────────────────────────────────────────────────────────────┐
│                      USER INTERFACE LAYER                     │
├────────────────┬─────────────────┬───────────────────────────┤
│   Login Dialog │  Main Window    │  Analytics Dashboard      │
│   ┌─────────┐  │  ┌───────────┐  │  ┌──────────────────┐    │
│   │Designer │  │  │ Toolbar   │  │  │ 10 Chart Types   │    │
│   │  User   │  │  │ Canvas    │  │  │ 6 KPI Cards      │    │
│   └─────────┘  │  │ Component │  │  │ Controls         │    │
│                │  │ Panel     │  │  └──────────────────┘    │
│                │  └───────────┘  │                           │
└────────────────┴─────────────────┴───────────────────────────┘
                           │
┌──────────────────────────┴───────────────────────────────────┐
│                     BUSINESS LOGIC LAYER                      │
├─────────────────────┬──────────────────┬─────────────────────┤
│ Component Registry  │  Message Server  │  Analytics Engine   │
│ ┌─────────────────┐ │ ┌──────────────┐ │ ┌─────────────────┐ │
│ │ Load JSON       │ │ │ TCP Listener │ │ │ Health Tracking │ │
│ │ Validate Types  │ │ │ UDP Listener │ │ │ KPI Calculation │ │
│ │ Type Management │ │ │ JSON Parser  │ │ │ Chart Data Gen  │ │
│ └─────────────────┘ │ └──────────────┘ │ └─────────────────┘ │
├─────────────────────┴──────────────────┴─────────────────────┤
│        Canvas Manager    │    Theme Manager    │ Voice Alerts │
└──────────────────────────┴─────────────────────┴──────────────┘
                           │
┌──────────────────────────┴───────────────────────────────────┐
│                        DATA LAYER                             │
├───────────────────┬──────────────────┬───────────────────────┤
│ components.json   │  Design Files    │  Health Data Streams  │
│ ┌───────────────┐ │ ┌──────────────┐ │ ┌───────────────────┐ │
│ │ Component     │ │ │ .design      │ │ │ Real-time JSON    │ │
│ │ Definitions   │ │ │ Format       │ │ │ TCP/UDP Messages  │ │
│ └───────────────┘ │ └──────────────┘ │ └───────────────────┘ │
└───────────────────┴──────────────────┴───────────────────────┘
                           │
┌──────────────────────────┴───────────────────────────────────┐
│                    EXTERNAL SYSTEMS                           │
│  ┌──────────────┐  ┌──────────────┐  ┌──────────────┐       │
│  │  Python      │  │  Python      │  │  Python      │       │
│  │  Simulator 1 │  │  Simulator 2 │  │  Simulator N │       │
│  └──────────────┘  └──────────────┘  └──────────────┘       │
└──────────────────────────────────────────────────────────────┘
"""
    add_diagram_slide(prs, "Block Diagram", block_diagram)
    
    # Slide 15: UML Class Diagram - Core Classes
    uml_diagram1 = """
┌─────────────────────────────┐
│   ComponentRegistry         │ <<Singleton>>
├─────────────────────────────┤
│ - instance: static          │
│ - components: QVector       │
│ - componentsFile: QString   │
├─────────────────────────────┤
│ + getInstance(): static     │
│ + loadFromJson(): bool      │
│ + addComponentType()        │
│ + getDefinition(): Def      │
│ + getAllTypes(): QStringList│
└──────────┬──────────────────┘
           │ manages
           ▼
┌─────────────────────────────┐
│  ComponentDefinition        │ <<Struct>>
├─────────────────────────────┤
│ + type_id: QString          │
│ + display_name: QString     │
│ + label: QString            │
│ + description: QString      │
│ + image_dir: QString        │
│ + icon_color: QString       │
│ + subsystems: QStringList   │
│ + protocol: QString         │
│ + port: int                 │
│ + category: QString         │
│ + shape: QString            │
└─────────────────────────────┘

┌─────────────────────────────┐        ┌──────────────────────┐
│      MainWindow             │        │     Canvas           │
├─────────────────────────────┤        ├──────────────────────┤
│ - canvas: Canvas*           │◆──────│ - components: QList  │
│ - dashboard: Analytics*     │        │ - connections: QList │
│ - messageServer: Server*    │        ├──────────────────────┤
│ - currentRole: UserRole     │        │ + addComponent()     │
├─────────────────────────────┤        │ + drawConnection()   │
│ + showAnalyticsDashboard()  │        │ + saveToFile()       │
│ + toggleTheme()             │        │ + loadFromFile()     │
│ + switchMode()              │        └──────────────────────┘
└─────────────────────────────┘
"""
    add_diagram_slide(prs, "UML Class Diagram - Core Classes", uml_diagram1)
    
    # Slide 16: UML Class Diagram - Component Hierarchy
    uml_diagram2 = """
┌──────────────────────────┐
│   QGraphicsItem          │ <<Qt>>
└────────────┬─────────────┘
             │ inherits
             ▼
┌──────────────────────────┐
│      Component           │
├──────────────────────────┤
│ - id: QString            │
│ - type: QString          │
│ - color: QColor          │
│ - size: double           │
│ - subcomponents: QList   │
├──────────────────────────┤
│ + updateHealth(data)     │
│ + paint()                │
│ + boundingRect()         │
│ + mousePressEvent()      │
└────────┬─────────────────┘
         │ contains
         ▼
┌──────────────────────────┐
│     SubComponent         │
├──────────────────────────┤
│ - name: QString          │
│ - health: double         │
│ - status: QString        │
├──────────────────────────┤
│ + updateStatus()         │
│ + getHealth(): double    │
└──────────────────────────┘

┌──────────────────────────┐        ┌──────────────────────┐
│   MessageServer          │        │  AnalyticsDashboard  │
├──────────────────────────┤        ├──────────────────────┤
│ - tcpServer: QTcpServer* │        │ - charts: QChartView*│
│ - udpSocket: QUdpSocket* │        │ - kpiLabels: QLabel* │
│ - connections: QList     │        │ - healthData: QVector│
├──────────────────────────┤        ├──────────────────────┤
│ + startServer()          │        │ + updateCharts()     │
│ + onTcpMessage()         │        │ + createHealthTrend()│
│ + onUdpMessage()         │        │ + createPieChart()   │
│ + processHealthData()    │        │ + calculateKPIs()    │
└──────────────────────────┘        └──────────────────────┘
"""
    add_diagram_slide(prs, "UML Class Diagram - Component Hierarchy", uml_diagram2)
    
    # Slide 17: Sequence Diagram - Component Addition
    seq_diagram1 = """
Designer     ComponentList   Canvas    ComponentRegistry
   │               │            │              │
   │ Click "+Add"  │            │              │
   ├──────────────►│            │              │
   │               │ Show Dialog│              │
   │               ├───────────►│              │
   │               │            │              │
   │ Fill Form     │            │              │
   │ Click Add     │            │              │
   │               │            │  addType()   │
   │               │            ├─────────────►│
   │               │            │              │
   │               │            │  Validate    │
   │               │            │  Save JSON   │
   │               │            │◄─────────────┤
   │               │            │  Success     │
   │               │ refreshList│              │
   │               │◄───────────┤              │
   │               │            │              │
   │ Drag New Type │            │              │
   ├──────────────►│            │              │
   │               │ createComp │              │
   │               ├───────────►│              │
   │               │            │ getDefinition│
   │               │            ├─────────────►│
   │               │            │◄─────────────┤
   │               │            │ Returns Data │
   │               │◄───────────┤              │
   │               │ Component  │              │
   │◄──────────────┤ Added      │              │
   │ Visual Update │            │              │
"""
    add_diagram_slide(prs, "Sequence Diagram - Component Addition", seq_diagram1)
    
    # Slide 18: Sequence Diagram - Health Update
    seq_diagram2 = """
External      MessageServer   Canvas    Component   Analytics
System            │            │           │            │
  │               │            │           │            │
  │ Send JSON     │            │           │            │
  ├──────────────►│            │           │            │
  │   (TCP/UDP)   │            │           │            │
  │               │ Parse JSON │           │            │
  │               │ Validate   │           │            │
  │               │            │           │            │
  │               │ Find Comp  │           │            │
  │               ├───────────►│           │            │
  │               │            │ getByID() │            │
  │               │            ├──────────►│            │
  │               │            │◄──────────┤            │
  │               │            │  Return   │            │
  │               │◄───────────┤           │            │
  │               │            │           │            │
  │               │ updateHealth()         │            │
  │               ├────────────────────────►│            │
  │               │            │           │ SetColor   │
  │               │            │           │ SetSize    │
  │               │            │           │ Repaint    │
  │               │            │           │            │
  │               │ recordEvent()          │            │
  │               ├────────────────────────────────────►│
  │               │            │           │            │
  │               │            │           │ UpdateKPIs │
  │               │            │           │ UpdateChart│
  │◄──────────────┤            │           │            │
  │   ACK         │            │           │            │
"""
    add_diagram_slide(prs, "Sequence Diagram - Health Update", seq_diagram2)
    
    # Slide 19: Sequence Diagram - Dashboard View
    seq_diagram3 = """
User      MainWindow   AnalyticsDashboard   Charts   DataEngine
 │            │                │               │          │
 │ Click      │                │               │          │
 │ Dashboard  │                │               │          │
 ├───────────►│                │               │          │
 │            │ Show Dashboard │               │          │
 │            ├───────────────►│               │          │
 │            │                │ Initialize    │          │
 │            │                │               │          │
 │            │                │ Create Charts │          │
 │            │                ├──────────────►│          │
 │            │                │               │ Generate │
 │            │                │               │ SampleData
 │            │                │               ├─────────►│
 │            │                │               │◄─────────┤
 │            │                │◄──────────────┤ Data     │
 │            │                │ Charts Ready  │          │
 │            │                │               │          │
 │            │                │ Calculate KPIs│          │
 │            │                │               │          │
 │            │◄───────────────┤               │          │
 │◄───────────┤ Dashboard Open │               │          │
 │            │                │               │          │
 │            │  Timer (5sec)  │               │          │
 │            │                │◄───────────────          │
 │            │                │ Auto Refresh  │          │
 │            │                │ Update Charts │          │
 │            │                │ Update KPIs   │          │
"""
    add_diagram_slide(prs, "Sequence Diagram - Dashboard View", seq_diagram3)
    
    # Slide 20: Deployment Architecture
    deployment_diagram = """
┌────────────────────────────────────────────────────────────┐
│                   Deployment Environment                    │
│                                                             │
│  ┌────────────────────────────────────────────────────┐   │
│  │              Client Machine (Workstation)          │   │
│  │  ┌──────────────────────────────────────────────┐  │   │
│  │  │         Operating System (Linux/Windows)     │  │   │
│  │  │  ┌────────────────────────────────────────┐  │  │   │
│  │  │  │      Qt Runtime Environment            │  │  │   │
│  │  │  │  ┌──────────────────────────────────┐  │  │  │   │
│  │  │  │  │   UnifiedApp Executable          │  │  │  │   │
│  │  │  │  │                                  │  │  │  │   │
│  │  │  │  │  • GUI Components (Qt Widgets)  │  │  │  │   │
│  │  │  │  │  • Network Listeners (Qt Net)   │  │  │  │   │
│  │  │  │  │  • Charts (Qt Charts)           │  │  │  │   │
│  │  │  │  │                                  │  │  │  │   │
│  │  │  │  │  Binds to:                       │  │  │  │   │
│  │  │  │  │  • TCP Port 12345                │  │  │  │   │
│  │  │  │  │  • UDP Port 12346                │  │  │  │   │
│  │  │  │  └──────────────┬───────────────────┘  │  │  │   │
│  │  │  └────────────────┼────────────────────────┘  │  │   │
│  │  │                   │                           │  │   │
│  │  │  ┌────────────────┼────────────────────────┐  │  │   │
│  │  │  │  Configuration Files                    │  │  │   │
│  │  │  │  • components.json                      │  │  │   │
│  │  │  │  • *.design files                       │  │  │   │
│  │  │  │  • styles.qss                           │  │  │   │
│  │  │  └─────────────────────────────────────────┘  │  │   │
│  │  └──────────────────────────────────────────────┘  │   │
│  └────────────────────────────────────────────────────┘   │
│                                                             │
│  ┌────────────────────────────────────────────────────┐   │
│  │        External Simulator Machine (Optional)       │   │
│  │  ┌──────────────────────────────────────────────┐  │   │
│  │  │         Python 3.x Runtime                   │  │   │
│  │  │  ┌────────────────────────────────────────┐  │  │   │
│  │  │  │  Simulator Scripts                     │  │  │   │
│  │  │  │  • external_system.py                  │  │  │   │
│  │  │  │  • run_multiple_systems.py             │  │  │   │
│  │  │  │  • apcu_simulator.py                   │  │  │   │
│  │  │  │                                        │  │  │   │
│  │  │  │  Connects to:                          │  │  │   │
│  │  │  │  • localhost:12345 (TCP)               │  │  │   │
│  │  │  │  • localhost:12346 (UDP)               │  │  │   │
│  │  │  └────────────┬───────────────────────────┘  │  │   │
│  │  └───────────────┼──────────────────────────────┘  │   │
│  └──────────────────┼─────────────────────────────────┘   │
│                     │                                      │
│  ┌──────────────────┼─────────────────────────────────┐   │
│  │  Network Layer   │                                 │   │
│  │  ┌───────────────▼──────────────────────────────┐  │   │
│  │  │     TCP/IP Stack (localhost or network)     │  │   │
│  │  │     • TCP Sockets (reliable delivery)       │  │   │
│  │  │     • UDP Sockets (low latency)             │  │   │
│  │  └─────────────────────────────────────────────┘  │   │
│  └────────────────────────────────────────────────────┘   │
└────────────────────────────────────────────────────────────┘

Network Topology:
┌─────────────┐         TCP/UDP          ┌──────────────┐
│  UnifiedApp │◄──────────────────────────│  Simulator 1 │
│  (Server)   │                           └──────────────┘
│             │         TCP/UDP          ┌──────────────┐
│ Port 12345  │◄──────────────────────────│  Simulator 2 │
│ Port 12346  │                           └──────────────┘
│             │         TCP/UDP          ┌──────────────┐
│             │◄──────────────────────────│  Simulator N │
└─────────────┘                           └──────────────┘
"""
    add_diagram_slide(prs, "Deployment Architecture", deployment_diagram)
    
    # Slide 21: Technology Stack
    add_two_column_slide(prs, "Technology Stack",
        [
            "Frontend / UI Framework:",
            "• Qt 5.x / 6.x",
            "  - Qt Widgets (GUI)",
            "  - Qt Network (TCP/UDP)",
            "  - Qt Charts (Visualization)",
            "",
            "Programming Languages:",
            "• C++ (Core Application)",
            "  - C++11/14/17 standards",
            "  - Object-oriented design",
            "",
            "• Python 3.x (Simulators)",
            "  - Standard library only",
            "  - Socket programming",
            "",
            "Data Formats:",
            "• JSON (Configuration & Messages)",
            "• Custom .design format"
        ],
        [
            "Build System:",
            "• qmake (Qt project files)",
            "• make / nmake",
            "",
            "Key Libraries:",
            "• Qt Core (Base functionality)",
            "• Qt GUI (Graphics)",
            "• Qt Widgets (UI components)",
            "• Qt Network (Networking)",
            "• Qt Charts (Data visualization)",
            "",
            "Development Tools:",
            "• Qt Creator (IDE)",
            "• g++ / clang / MSVC (Compilers)",
            "• Git (Version control)",
            "",
            "Testing:",
            "• Shell scripts for integration tests",
            "• Python simulators for health data"
        ]
    )
    
    # Slide 22: Component Registry JSON Schema
    json_schema = """
{
  "version": "2.0",
  "description": "Modular component registry",
  "components": [
    {
      "type_id": "Antenna",                // Unique identifier
      "display_name": "Antenna System",    // Human-readable name
      "label": "ANT",                      // Short canvas label
      "description": "Phased array...",    // Description
      "image_dir": "antenna",              // Image directory
      "icon_color": "#4CAF50",             // Fallback color
      "subsystems": [                      // Sub-components
        "APCU Controller",
        "Temperature Monitoring",
        "Board Data"
      ],
      "allowed_widgets": [                 // Design widgets
        "Label", "LineEdit", "Button"
      ],
      "protocol": "TCP",                   // Health protocol
      "port": 12345,                       // Port number
      "category": "Sensor",                // Category
      "shape": "ellipse"                   // Fallback shape
    }
  ]
}

Benefits:
• Zero-code component addition
• Data-driven rendering
• Backward compatibility
• Extensible protocol support
"""
    add_diagram_slide(prs, "Component Registry JSON Schema", json_schema)
    
    # Slide 23: Health Status Protocol
    add_content_slide(prs, "Health Status Protocol", [
        "• JSON-based message format",
        "• Message structure:",
        {"text": "{ \"component_id\": \"comp_1\", \"color\": \"#00FF00\", \"size\": 95.5 }", "level": 1},
        "• Health level mapping:",
        {"text": "90-100%: Operational (Green #00FF00)", "level": 1},
        {"text": "70-89%: Warning (Yellow #FFFF00)", "level": 1},
        {"text": "40-69%: Degraded (Orange #FFA500)", "level": 1},
        {"text": "10-39%: Critical (Red #FF0000)", "level": 1},
        {"text": "0-9%: Offline (Gray #808080)", "level": 1},
        "• Protocol options:",
        {"text": "TCP (Port 12345): Reliable, line-delimited", "level": 1},
        {"text": "UDP (Port 12346): Low-latency, fire-and-forget", "level": 1},
        "• Multi-client support with concurrent connections"
    ])
    
    # Slide 24: Project Structure
    project_structure = """
.
├── UnifiedApp/                    # Main Qt Application
│   ├── UnifiedApp.pro             # Qt project file
│   ├── components.json            # Component registry
│   ├── main.cpp                   # Application entry point
│   ├── componentregistry.h/cpp    # Singleton registry
│   ├── mainwindow.h/cpp           # Main window
│   ├── component.h/cpp            # Component graphics item
│   ├── canvas.h/cpp               # Designer canvas
│   ├── messageserver.h/cpp        # Multi-protocol server
│   ├── analyticsdashboard.h/cpp   # Analytics dashboard
│   ├── logindialog.h/cpp          # Login with roles
│   ├── thememanager.h/cpp         # Theme switching
│   ├── subcomponent.h/cpp         # Sub-component support
│   ├── connection.h/cpp           # Component connections
│   ├── analytics.h/cpp            # Analytics widget
│   └── styles*.qss                # Theme stylesheets
│
├── ExternalSystems/               # Python simulators
│   ├── external_system.py         # Single component simulator
│   ├── run_multiple_systems.py    # Multi-component launcher
│   └── apcu_simulator.py          # APCU-specific simulator
│
├── assets/subsystems/             # Component images
│   ├── antenna/
│   ├── power_system/
│   ├── liquid_cooling_unit/
│   ├── communication_system/
│   └── radar_computer/
│
└── Documentation/                 # Project docs
    ├── README.md
    ├── IMPLEMENTATION_SUMMARY.md
    ├── ANALYTICS_DASHBOARD_GUIDE.md
    └── DASHBOARD_FEATURES_SUMMARY.md
"""
    add_diagram_slide(prs, "Project Structure", project_structure)
    
    # Slide 25: Build and Deployment
    add_content_slide(prs, "Build and Deployment", [
        "• Prerequisites:",
        {"text": "Qt 5.x or 6.x with development tools", "level": 1},
        {"text": "C++ compiler (g++, clang, or MSVC)", "level": 1},
        {"text": "Python 3.x for simulators", "level": 1},
        "• Build steps:",
        {"text": "1. cd UnifiedApp", "level": 1},
        {"text": "2. qmake UnifiedApp.pro", "level": 1},
        {"text": "3. make (or nmake on Windows)", "level": 1},
        {"text": "4. ./UnifiedApp", "level": 1},
        "• Installation packages:",
        {"text": "Ubuntu: sudo apt-get install qt5-default qtcharts5-dev", "level": 1},
        {"text": "Fedora: sudo dnf install qt5-qtbase-devel", "level": 1},
        {"text": "macOS: brew install qt@5", "level": 1}
    ])
    
    # Slide 26: Usage Workflow
    add_content_slide(prs, "Complete Usage Workflow", [
        "• Step 1: Launch Application",
        {"text": "./UnifiedApp", "level": 1},
        "• Step 2: Login",
        {"text": "Designer (Designer/designer) for design mode", "level": 1},
        {"text": "User (User/user) for runtime monitoring", "level": 1},
        "• Step 3: Design System (Designer mode)",
        {"text": "Drag components from panel to canvas", "level": 1},
        {"text": "Draw connections between components", "level": 1},
        {"text": "Add custom component types if needed", "level": 1},
        {"text": "Save design file (.design)", "level": 1},
        "• Step 4: Monitor System (Runtime mode)",
        {"text": "Load saved design", "level": 1},
        {"text": "Server auto-starts on ports 12345/12346", "level": 1},
        "• Step 5: Start Simulators",
        {"text": "python3 run_multiple_systems.py --components comp_1 comp_2", "level": 1}
    ])
    
    # Slide 27: Key Benefits
    add_content_slide(prs, "Key Benefits", [
        "• Modularity",
        {"text": "Add components without code changes", "level": 1},
        {"text": "JSON-based configuration", "level": 1},
        "• Flexibility",
        {"text": "Multi-protocol support (TCP, UDP, WebSocket, MQTT)", "level": 1},
        {"text": "Customizable component types and shapes", "level": 1},
        "• Real-time Monitoring",
        {"text": "Live health status updates", "level": 1},
        {"text": "Visual color-coded indicators", "level": 1},
        "• Rich Analytics",
        {"text": "10 different chart types", "level": 1},
        {"text": "Comprehensive KPI tracking", "level": 1},
        "• User Experience",
        {"text": "Intuitive drag-and-drop designer", "level": 1},
        {"text": "Dark/Light theme support", "level": 1}
    ])
    
    # Slide 28: Performance Considerations
    add_content_slide(prs, "Performance Considerations", [
        "• Efficient data structures",
        {"text": "QVector for component storage", "level": 1},
        {"text": "Hash-based component lookup", "level": 1},
        "• Optimized rendering",
        {"text": "Qt Graphics View framework", "level": 1},
        {"text": "Smart repaint regions", "level": 1},
        "• Network optimization",
        {"text": "Asynchronous message processing", "level": 1},
        {"text": "Multi-threaded server connections", "level": 1},
        "• Memory management",
        {"text": "Qt parent-child ownership", "level": 1},
        {"text": "Smart pointers for dynamic objects", "level": 1},
        "• Chart rendering",
        {"text": "Data sampling for large datasets", "level": 1},
        {"text": "Optimized refresh intervals (5 seconds)", "level": 1}
    ])
    
    # Slide 29: Security Considerations
    add_content_slide(prs, "Security Considerations", [
        "• Authentication",
        {"text": "Role-based access control (Designer/User)", "level": 1},
        {"text": "Login dialog with credentials", "level": 1},
        "• Input validation",
        {"text": "JSON schema validation", "level": 1},
        {"text": "Component ID verification", "level": 1},
        "• Network security",
        {"text": "Localhost by default (can be configured)", "level": 1},
        {"text": "Protocol validation for health messages", "level": 1},
        "• Configuration security",
        {"text": "Restricted file access permissions", "level": 1},
        {"text": "Design file validation on load", "level": 1},
        "• Future enhancements:",
        {"text": "TLS/SSL for network communication", "level": 1},
        {"text": "User database with hashed passwords", "level": 1},
        {"text": "Audit logging for security events", "level": 1}
    ])
    
    # Slide 30: Future Enhancements
    add_content_slide(prs, "Future Enhancements", [
        "• Real-time data integration",
        {"text": "Connect dashboard to live component data", "level": 1},
        {"text": "Historical data persistence (database)", "level": 1},
        "• Advanced analytics",
        {"text": "Predictive health analytics with ML", "level": 1},
        {"text": "Anomaly detection algorithms", "level": 1},
        "• Enhanced protocols",
        {"text": "WebSocket implementation", "level": 1},
        {"text": "MQTT broker integration", "level": 1},
        "• Export and reporting",
        {"text": "CSV/Excel data export", "level": 1},
        {"text": "PDF report generation with charts", "level": 1},
        "• Customization",
        {"text": "User-configurable dashboard layouts", "level": 1},
        {"text": "Custom chart creation wizard", "level": 1}
    ])
    
    # Slide 31: Testing Strategy
    add_content_slide(prs, "Testing Strategy", [
        "• Unit Testing",
        {"text": "Component Registry validation", "level": 1},
        {"text": "JSON parsing and validation", "level": 1},
        "• Integration Testing",
        {"text": "Message server with simulators", "level": 1},
        {"text": "Component health update flow", "level": 1},
        {"text": "Shell scripts: test_component_hierarchy.sh", "level": 1},
        "• System Testing",
        {"text": "End-to-end workflow validation", "level": 1},
        {"text": "Designer to Runtime workflow", "level": 1},
        "• UI Testing",
        {"text": "Theme switching validation", "level": 1},
        {"text": "Dashboard rendering tests", "level": 1},
        "• Performance Testing",
        {"text": "Multi-client connection load", "level": 1},
        {"text": "Chart rendering with large datasets", "level": 1}
    ])
    
    # Slide 32: Documentation
    add_content_slide(prs, "Project Documentation", [
        "• README.md",
        {"text": "Comprehensive project overview", "level": 1},
        {"text": "Quick start guide", "level": 1},
        {"text": "Feature descriptions", "level": 1},
        "• IMPLEMENTATION_SUMMARY.md",
        {"text": "Detailed implementation notes", "level": 1},
        {"text": "Code statistics and metrics", "level": 1},
        "• ANALYTICS_DASHBOARD_GUIDE.md",
        {"text": "Dashboard usage and features", "level": 1},
        {"text": "Chart type descriptions", "level": 1},
        "• DASHBOARD_FEATURES_SUMMARY.md",
        {"text": "Quick reference for dashboard", "level": 1},
        "• Code Comments",
        {"text": "Inline documentation in source files", "level": 1},
        {"text": "Function and class descriptions", "level": 1}
    ])
    
    # Slide 33: Conclusion
    add_content_slide(prs, "Conclusion", [
        "• Comprehensive radar system monitoring solution",
        "• Modern architecture with proven design patterns",
        "• Modular and extensible by design",
        "• Rich real-time visualization capabilities",
        "• Multi-protocol support for flexibility",
        "• Professional UI/UX with theme support",
        "• Well-documented and maintainable codebase",
        "",
        "Key Achievements:",
        {"text": "✓ Zero-code component extension via JSON", "level": 1},
        {"text": "✓ Real-time multi-protocol health monitoring", "level": 1},
        {"text": "✓ Advanced analytics with 10 chart types", "level": 1},
        {"text": "✓ Designer and Runtime modes in unified app", "level": 1},
        {"text": "✓ Comprehensive documentation and testing", "level": 1}
    ])
    
    # Slide 34: Thank You
    add_title_slide(
        prs,
        "Thank You",
        "Radar System Monitoring Application\n\nQuestions & Discussion"
    )
    
    return prs

def main():
    """Main function to create and save the presentation"""
    print("Creating PowerPoint presentation...")
    
    prs = create_presentation()
    
    output_file = "/workspace/Radar_System_Architecture_Presentation.pptx"
    prs.save(output_file)
    
    print(f"✓ Presentation created successfully!")
    print(f"✓ Saved to: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
